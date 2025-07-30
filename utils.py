from docx import Document
from pptx import Presentation
import csv
import fitz  # PyMuPDF

# === File Parser ===
class FileParser:
    def __init__(self):
        """Initialize the file parser for various document formats."""
        print("[INIT] FileParser initialized.")

    def parse_pdf(self, file_path, skip_start_pages=0, skip_end_pages=0):
        """
        Extract text from PDF file with option to skip pages from start and end.
        
        Args:
            file_path (str): Path to the PDF file
            skip_start_pages (int): Number of pages to skip from the beginning
            skip_end_pages (int): Number of pages to skip from the end
        
        Returns:
            str: Extracted text content
        """
        doc = fitz.open(file_path)
        total_pages = len(doc)
        
        # Calculate page range
        start_page = skip_start_pages
        end_page = total_pages - skip_end_pages
        
        if start_page >= end_page:
            raise ValueError(f"Invalid page range: start_page={start_page}, end_page={end_page}")
        
        pages_text = []
        for page_num in range(start_page, end_page):
            pages_text.append(doc[page_num].get_text())
        
        doc.close()
        print(pages_text[:20])
        return "\n".join(pages_text)

    def parse_docx(self, file_path, skip_start_pages=0, skip_end_pages=0):
        """
        Extract text from Word document (.docx) with option to skip pages.
        Note: DOCX doesn't have explicit pages, so we treat paragraphs as units.
        
        Args:
            file_path (str): Path to the DOCX file
            skip_start_pages (int): Number of paragraphs to skip from the beginning
            skip_end_pages (int): Number of paragraphs to skip from the end
        
        Returns:
            str: Extracted text content
        """
        doc = Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        total_paragraphs = len(paragraphs)
        
        # Calculate paragraph range
        start_para = skip_start_pages
        end_para = total_paragraphs - skip_end_pages
        
        if start_para >= end_para:
            print(f"[WARNING] Invalid paragraph range for {file_path}: start={start_para}, end={end_para}")
            return ""
        
        selected_paragraphs = paragraphs[start_para:end_para]
        return "\n".join(selected_paragraphs)

    def parse_pptx(self, file_path, skip_start_pages=0, skip_end_pages=0):
        """
        Extract text from PowerPoint presentation (.pptx) with option to skip slides.
        
        Args:
            file_path (str): Path to the PPTX file
            skip_start_pages (int): Number of slides to skip from the beginning
            skip_end_pages (int): Number of slides to skip from the end
        
        Returns:
            str: Extracted text content from selected slides
        """
        prs = Presentation(file_path)
        slides = list(prs.slides)
        total_slides = len(slides)
        
        # Calculate slide range
        start_slide = skip_start_pages
        end_slide = total_slides - skip_end_pages
        
        if start_slide >= end_slide:
            print(f"[WARNING] Invalid slide range for {file_path}: start={start_slide}, end={end_slide}")
            return ""
        
        text = ""
        for slide_idx in range(start_slide, end_slide):
            slide = slides[slide_idx]
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def parse_csv(self, file_path, skip_start_pages=0, skip_end_pages=0):
        """
        Extract text from CSV file with option to skip rows from start and end.
        
        Args:
            file_path (str): Path to the CSV file
            skip_start_pages (int): Number of rows to skip from the beginning
            skip_end_pages (int): Number of rows to skip fr
            om the end
        
        Returns:
            str: Extracted text content with selected rows separated by newlines
        """
        rows = []
        with open(file_path, newline='', encoding='utf-8') as csvfile:
            reader = csv.reader(csvfile)
            all_rows = list(reader)
        
        total_rows = len(all_rows)
        
        # Calculate row range
        start_row = skip_start_pages
        end_row = total_rows - skip_end_pages
        
        if start_row >= end_row:
            print(f"[WARNING] Invalid row range for {file_path}: start={start_row}, end={end_row}")
            return ""
        
        selected_rows = all_rows[start_row:end_row]
        processed_rows = [" ".join(row) for row in selected_rows]
        return "\n".join(processed_rows)

    def parse_txt(self, file_path, skip_start_pages=0, skip_end_pages=0):
        """
        Extract text from plain text file with option to skip lines from start and end.
        
        Args:
            file_path (str): Path to the text file
            skip_start_pages (int): Number of lines to skip from the beginning
            skip_end_pages (int): Number of lines to skip from the end
        
        Returns:
            str: Selected lines joined as string
        """
        with open(file_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        
        total_lines = len(lines)
        
        # Calculate line range
        start_line = skip_start_pages
        end_line = total_lines - skip_end_pages
        
        if start_line >= end_line:
            print(f"[WARNING] Invalid line range for {file_path}: start={start_line}, end={end_line}")
            return ""
        
        selected_lines = lines[start_line:end_line]
        return "".join(selected_lines)

    def parse(self, file_path, skip_start_pages=0, skip_end_pages=0):
        """
        Parse file based on its extension and extract text content.
        
        Args:
            file_path (str): Path to the file
            skip_start_pages (int): Number of units to skip from start (pages/paragraphs/slides/rows/lines)
            skip_end_pages (int): Number of units to skip from end (pages/paragraphs/slides/rows/lines)
        
        Returns:
            str: Extracted text content
        
        Raises:
            ValueError: If file type is not supported
        """
        ext = file_path.lower().split('.')[-1]
        if ext == 'pdf':
            return self.parse_pdf(file_path, skip_start_pages, skip_end_pages)
        elif ext == 'docx':
            return self.parse_docx(file_path, skip_start_pages, skip_end_pages)
        elif ext == 'pptx':
            return self.parse_pptx(file_path, skip_start_pages, skip_end_pages)
        elif ext == 'csv':
            return self.parse_csv(file_path, skip_start_pages, skip_end_pages)
        elif ext in ['txt', 'md']:
            return self.parse_txt(file_path, skip_start_pages, skip_end_pages)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
        

# t=FileParser()
# text=(t.parse("docs\9241544228_eng.pdf",skip_end_pages=3,skip_start_pages=0))      